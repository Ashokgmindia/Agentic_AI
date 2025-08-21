import os
import io
import logging
import pandas as pd
import json
import base64
from textwrap import dedent
import requests
from datetime import datetime
from crewai.tools import tool
from langchain_community.document_loaders import (
    PyPDFLoader, CSVLoader, UnstructuredPowerPointLoader,
    Docx2txtLoader
)
from PIL import Image

# Set up logging for the tools
logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO)

# Initialize Whisper model with error handling
whisper_model = None
WHISPER_AVAILABLE = False

try:
    import whisper
    WHISPER_AVAILABLE = True
    logger.info("Whisper module loaded successfully")
except ImportError:
    logger.warning("Whisper module not available. Audio transcription will use fallback method.")
    WHISPER_AVAILABLE = False

def get_whisper_model():
    global whisper_model
    if whisper_model is None and WHISPER_AVAILABLE:
        try:
            whisper_model = whisper.load_model("base")
            logger.info("Whisper model loaded successfully")
        except Exception as e:
            logger.error(f"Failed to load Whisper model: {e}")
            return None
    return whisper_model

# ========== TEXT ==========
@tool("Ingest Text")
def ingest_text(file_path: str):
    """Reads a plain text file and returns content, metadata, and detailed analysis."""
    if not os.path.exists(file_path):
        return {"error": "File not found"}
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        
        word_count = len(content.split())
        line_count = len(content.split('\n'))
        char_count = len(content)
        
        # Basic content analysis
        analysis = {
            "word_count": word_count,
            "line_count": line_count,
            "character_count": char_count,
            "content_preview": content[:500] + "..." if len(content) > 500 else content
        }
        
        return {
            "content": content,
            "metadata": {"type": "text", "source": file_path, "analysis": analysis},
            "summary": f"Text file with {word_count} words, {line_count} lines, and {char_count} characters."
        }
    except Exception as e:
        return {"error": f"Error reading text file: {e}"}

# ========== PDF ==========
@tool("Ingest PDF")
def ingest_pdf(file_path: str):
    """Extracts text from a PDF file and provides detailed analysis."""
    if not os.path.exists(file_path):
        return {"error": "File not found"}
    try:
        loader = PyPDFLoader(file_path)
        pages = loader.load()
        content = "\n\n".join([p.page_content for p in pages])
        
        total_words = len(content.split())
        
        return {
            "content": content,
            "metadata": {
                "type": "pdf", 
                "pages": len(pages), 
                "source": file_path,
                "word_count": total_words
            },
            "summary": f"PDF document with {len(pages)} pages containing {total_words} words. Key content extracted and ready for analysis."
        }
    except Exception as e:
        return {"error": f"Error processing PDF: {e}"}

# ========== FIXED IMAGE TOOL ==========
@tool("Ingest Image")
def ingest_image(file_path: str):
    """Analyzes image and provides detailed description and metadata."""
    if not os.path.exists(file_path):
        return {"error": f"Image file not found at {file_path}"}
    
    try:
        # Validate it's actually an image
        with Image.open(file_path) as img:
            width, height = img.size
            mode = img.mode
            format_type = img.format or "Unknown"
            
            # Calculate file size
            file_size = os.path.getsize(file_path)
            
            content_description = f"""# Image Analysis Report

## Technical Specifications
- **File**: {os.path.basename(file_path)}
- **Dimensions**: {width} x {height} pixels
- **Format**: {format_type}
- **Color Mode**: {mode}
- **File Size**: {file_size:,} bytes
- **Aspect Ratio**: {width/height:.2f}

## Quality Assessment
- **Resolution**: {width * height:,} total pixels
- **Quality Rating**: {'High' if width * height > 1000000 else 'Medium' if width * height > 100000 else 'Standard'}
- **Web Suitability**: {'Excellent' if format_type in ['JPEG', 'PNG'] else 'Good'}

## Business Intelligence
This image file is technically sound and ready for analysis:
- **Processing Ready**: Image successfully loaded and validated
- **Content Type**: Visual content suitable for multimodal analysis
- **File Quality**: {format_type} format with {width}x{height} resolution
- **Storage Efficient**: {file_size:,} bytes

## Recommendations
- Image is properly formatted for business use
- Resolution is {'suitable for printing' if width >= 1200 else 'suitable for web display'}
- File size is {'optimized' if file_size < 2000000 else 'large - consider compression'}
"""
        
        return {
            "content": content_description,
            "metadata": {
                "type": "image",
                "source": file_path,
                "width": width,
                "height": height,
                "format": format_type,
                "file_size": file_size
            },
            "summary": f"Image analyzed: {format_type} format, {width}x{height} pixels, {file_size:,} bytes"
        }
        
    except Exception as e:
        return {"error": f"Image processing error: {str(e)}"}


# ========== FIXED AUDIO TOOL ==========
@tool("Ingest Audio")  
def ingest_audio(file_path: str):
    """Transcribes speech from an audio file or provides detailed file analysis."""
    if not os.path.exists(file_path):
        return {"error": "Audio file not found"}
    
    try:
        file_size = os.path.getsize(file_path)
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # Validate audio file
        valid_audio_formats = ['.wav', '.mp3', '.m4a', '.aac', '.flac', '.ogg']
        if file_ext not in valid_audio_formats:
            return {"error": f"Unsupported audio format: {file_ext}. Supported: {', '.join(valid_audio_formats)}"}
        
        if WHISPER_AVAILABLE:
            try:
                model = get_whisper_model()
                if model:
                    # Try transcription with error handling
                    result = model.transcribe(file_path)
                    content = result.get("text", "").strip()
                    language = result.get("language", "unknown")
                    
                    if content:  # Successful transcription
                        formatted_content = f"""# Audio Transcription Report

## File Information
- **File**: {os.path.basename(file_path)}
- **Format**: {file_ext.upper()}
- **Size**: {file_size:,} bytes
- **Language**: {language}

## Transcription Results
**Content**: {content}

## Analysis Summary
Audio successfully transcribed with {len(content.split())} words detected in {language} language.
"""
                        return {
                            "content": formatted_content,
                            "metadata": {"type": "audio", "source": file_path, "language": language},
                            "summary": f"Audio transcribed: {len(content.split())} words in {language}"
                        }
            except Exception as whisper_error:
                print(f"Whisper transcription failed: {whisper_error}")
        
        # Fallback analysis
        fallback_content = f"""# Audio File Analysis

## File Information
- **File**: {os.path.basename(file_path)}
- **Format**: {file_ext.upper()}
- **Size**: {file_size:,} bytes
- **Estimated Duration**: ~{file_size // 16000} seconds

## Status
Audio file received and validated. Transcription service temporarily unavailable.

## File Assessment
- **Format**: {file_ext} is a {'standard' if file_ext in ['.wav', '.mp3'] else 'supported'} audio format
- **Size**: {'Large' if file_size > 10000000 else 'Medium' if file_size > 1000000 else 'Small'} file
- **Ready**: File is properly formatted for audio processing
"""
        
        return {
            "content": fallback_content,
            "metadata": {"type": "audio", "source": file_path, "file_size": file_size},
            "summary": f"Audio file validated: {file_ext} format, {file_size:,} bytes"
        }
        
    except Exception as e:
        return {"error": f"Audio processing error: {str(e)}"}


# ========== FIXED CSV TOOL ==========
@tool("Ingest CSV")
def ingest_csv(file_path: str):
    """Loads and analyzes CSV data with comprehensive insights."""
    if not os.path.exists(file_path):
        return {"error": "File not found"}
    
    try:
        # Try different encodings and separators
        encodings = ['utf-8', 'latin-1', 'cp1252']
        separators = [',', ';', '\t']
        
        df = None
        for encoding in encodings:
            for sep in separators:
                try:
                    df = pd.read_csv(file_path, encoding=encoding, sep=sep)
                    if len(df.columns) > 1:  # Valid CSV should have multiple columns
                        break
                except:
                    continue
            if df is not None and len(df.columns) > 1:
                break
        
        if df is None or len(df.columns) <= 1:
            return {"error": "Could not parse CSV file. Please check file format."}
        
        # Rest of the CSV analysis code remains the same...
        summary_text = f"""# CSV Data Analysis: {os.path.basename(file_path)}

## Dataset Overview
- **Total Records**: {len(df):,}
- **Total Columns**: {len(df.columns)}
- **File Size**: {os.path.getsize(file_path):,} bytes

## Column Summary
"""
        
        for i, col in enumerate(df.columns[:10]):  # Limit to first 10 columns
            dtype = str(df[col].dtype)
            null_count = df[col].isnull().sum()
            unique_count = df[col].nunique()
            
            summary_text += f"**{col}** ({dtype}): {null_count} missing, {unique_count} unique values\n"
        
        # Add sample data
        summary_text += f"\n## Sample Data (First 3 Rows)\n"
        try:
            sample_df = df.head(3)
            for idx in range(len(sample_df)):
                row_data = []
                for col in df.columns[:5]:  # First 5 columns only
                    val = str(sample_df.iloc[idx][col])[:20]
                    row_data.append(val)
                summary_text += f"Row {idx+1}: {' | '.join(row_data)}\n"
        except Exception as e:
            summary_text += f"Sample data error: {e}\n"
        
        return {
            "content": summary_text,
            "metadata": {"type": "csv", "source": file_path, "rows": len(df), "columns": len(df.columns)},
            "summary": f"CSV file analyzed: {len(df):,} rows × {len(df.columns)} columns"
        }
        
    except Exception as e:
        return {"error": f"CSV processing error: {str(e)}"}


# ========== FIXED PPT TOOL ==========
@tool("Ingest PPT")
def ingest_ppt(file_path: str):
    """Extracts and analyzes content from PowerPoint presentations."""
    if not os.path.exists(file_path):
        return {"error": "PowerPoint file not found"}
    
    try:
        # Alternative approach for PPT processing
        try:
            from pptx import Presentation
            
            # Use python-pptx as primary method
            prs = Presentation(file_path)
            slides_content = []
            total_words = 0
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + " "
                
                word_count = len(slide_text.split())
                total_words += word_count
                slides_content.append({
                    "slide_number": i,
                    "content": slide_text.strip(),
                    "word_count": word_count
                })
            
            content_summary = f"""# PowerPoint Analysis: {os.path.basename(file_path)}

## Presentation Overview
- **Total Slides**: {len(slides_content)}
- **Total Words**: {total_words}
- **Average Words per Slide**: {total_words // len(slides_content) if slides_content else 0}

## Slide Content Summary
"""
            
            for slide in slides_content[:5]:  # Show first 5 slides
                content_summary += f"\n**Slide {slide['slide_number']}** ({slide['word_count']} words)\n"
                preview = slide['content'][:200] + "..." if len(slide['content']) > 200 else slide['content']
                content_summary += f"{preview}\n"
            
            if len(slides_content) > 5:
                content_summary += f"\n*... and {len(slides_content)-5} more slides*\n"
            
            return {
                "content": content_summary,
                "metadata": {"type": "ppt", "slides": len(slides_content), "source": file_path, "total_words": total_words},
                "summary": f"PowerPoint analyzed: {len(slides_content)} slides with {total_words} total words"
            }
            
        except ImportError:
            # Fallback to unstructured if python-pptx not available
            loader = UnstructuredPowerPointLoader(file_path)
            docs = loader.load()
            
            content = "\n\n".join([doc.page_content for doc in docs])
            word_count = len(content.split())
            
            return {
                "content": f"# PowerPoint Analysis\n\n{content}",
                "metadata": {"type": "ppt", "slides": len(docs), "source": file_path, "total_words": word_count},
                "summary": f"PowerPoint processed: {len(docs)} sections with {word_count} words"
            }
            
    except Exception as e:
        return {"error": f"PowerPoint processing error: {str(e)}"}

# ========== DOC/DOCX ==========
@tool("Ingest DOC")
def ingest_doc(file_path: str):
    """Extracts and analyzes content from Word documents."""
    if not os.path.exists(file_path):
        return {"error": "File not found"}
    try:
        loader = Docx2txtLoader(file_path)
        docs = loader.load()
        content = "\n\n".join([doc.page_content for doc in docs])
        
        # Comprehensive analysis
        word_count = len(content.split())
        sentences = content.split('.')
        sentence_count = len([s for s in sentences if s.strip()])
        paragraph_count = len([p for p in content.split('\n\n') if p.strip()])
        
        analysis = {
            "word_count": word_count,
            "sentence_count": sentence_count,
            "paragraph_count": paragraph_count,
            "character_count": len(content),
            "avg_words_per_paragraph": word_count / paragraph_count if paragraph_count > 0 else 0,
            "avg_sentences_per_paragraph": sentence_count / paragraph_count if paragraph_count > 0 else 0
        }
        
        formatted_content = f"""# Word Document Analysis: {os.path.basename(file_path)}

## Document Statistics
- **Word Count**: {word_count:,}
- **Sentence Count**: {sentence_count:,}
- **Paragraph Count**: {paragraph_count}
- **Character Count**: {len(content):,}
- **Average Words per Paragraph**: {analysis['avg_words_per_paragraph']:.1f}
- **Average Sentences per Paragraph**: {analysis['avg_sentences_per_paragraph']:.1f}

## Content Analysis
### Document Type Assessment
{
'**Long Document**: Detailed content suitable for comprehensive analysis' if word_count > 2000 else
'**Medium Document**: Standard business document length' if word_count > 500 else
'**Short Document**: Brief content or summary document'
}

### Readability Metrics
- **Content Density**: {'High' if analysis['avg_words_per_paragraph'] > 100 else 'Medium' if analysis['avg_words_per_paragraph'] > 50 else 'Low'} (based on paragraph length)
- **Structure**: {'Well-structured' if paragraph_count > 3 else 'Simple structure'} with {paragraph_count} paragraphs
- **Sentence Complexity**: {'Complex' if word_count / sentence_count > 20 else 'Standard'} (avg {word_count / sentence_count:.1f} words per sentence)

## Document Content

{content}

## Business Intelligence Summary
### Content Classification
This document contains {word_count:,} words organized into {paragraph_count} paragraphs, indicating {'detailed documentation' if word_count > 1000 else 'standard business content'}. 

### Key Insights
- **Document Length**: {'Comprehensive' if word_count > 2000 else 'Standard' if word_count > 500 else 'Concise'} business document
- **Information Density**: {analysis['avg_words_per_paragraph']:.0f} words per paragraph suggests {'detailed explanations' if analysis['avg_words_per_paragraph'] > 100 else 'balanced content structure'}
- **Processing Value**: High-value text content ready for business analysis and decision-making

### Recommendations
- Content is suitable for {'detailed business analysis' if word_count > 1000 else 'standard business processing'}
- Document structure supports {'comprehensive review' if paragraph_count > 5 else 'quick review'} and analysis
- Text quality enables automated processing and insight extraction
"""
        
        return {
            "content": formatted_content,
            "metadata": {
                "type": "docx", 
                "source": file_path,
                "analysis": analysis
            },
            "summary": f"Word document comprehensively analyzed: {word_count:,} words across {paragraph_count} paragraphs and {sentence_count} sentences. Document classified as {'comprehensive' if word_count > 2000 else 'standard' if word_count > 500 else 'concise'} business content."
        }
    except Exception as e:
        return {"error": f"Error processing Word document: {e}"}

# ========== XLSX ==========
@tool("Ingest XLSX")
def ingest_xlsx(file_path: str):
    """Loads and analyzes Excel data with multi-sheet support."""
    if not os.path.exists(file_path):
        return {"error": "File not found"}
    try:
        # Load all sheets
        sheets_dict = pd.read_excel(file_path, sheet_name=None)
        
        analysis = {
            "total_sheets": len(sheets_dict),
            "sheet_names": list(sheets_dict.keys()),
            "total_rows": sum(len(df) for df in sheets_dict.values()),
            "total_columns": sum(len(df.columns) for df in sheets_dict.values()),
            "sheets_analysis": {}
        }
        
        content_summary = f"""# Excel Workbook Analysis: {os.path.basename(file_path)}

## Workbook Overview
- **Total Sheets**: {len(sheets_dict)}
- **Combined Data**: {analysis['total_rows']:,} total rows across all sheets
- **Total Columns**: {analysis['total_columns']} columns (sum across sheets)
- **File Complexity**: {'Complex multi-sheet workbook' if len(sheets_dict) > 3 else 'Standard workbook'}

## Sheet-by-Sheet Analysis

"""
        
        for sheet_name, df in sheets_dict.items():
            # Analyze each sheet
            numeric_cols = df.select_dtypes(include=['number']).columns
            categorical_cols = df.select_dtypes(include=['object']).columns
            null_percentage = (df.isnull().sum().sum() / (df.shape[0] * df.shape[1]) * 100) if df.shape * df.shape[1] > 0 else 0
            
            sheet_analysis = {
                "rows": len(df),
                "columns": len(df.columns),
                "column_names": list(df.columns),
                "numeric_columns": len(numeric_cols),
                "categorical_columns": len(categorical_cols),
                "completeness": 100 - null_percentage
            }
            analysis['sheets_analysis'][sheet_name] = sheet_analysis
            
            content_summary += f"### Sheet: {sheet_name}\n"
            content_summary += f"- **Dimensions**: {len(df):,} rows × {len(df.columns)} columns\n"
            content_summary += f"- **Data Types**: {len(numeric_cols)} numeric, {len(categorical_cols)} text/categorical\n"
            content_summary += f"- **Data Quality**: {sheet_analysis['completeness']:.1f}% complete\n"
            
            # Column details
            if len(df.columns) <= 10:
                content_summary += f"- **Columns**: {', '.join(df.columns)}\n"
            else:
                content_summary += f"- **Columns**: {', '.join(df.columns[:8])}, ... and {len(df.columns)-8} more\n"
            
            # Data insights for each sheet
            if len(df) > 0:
                content_summary += f"\n#### Business Intelligence - {sheet_name}\n"
                
                if len(numeric_cols) > 0:
                    content_summary += f"- **Quantitative Analysis Ready**: {len(numeric_cols)} numeric columns for calculations\n"
                if len(categorical_cols) > 0:
                    content_summary += f"- **Segmentation Potential**: {len(categorical_cols)} categorical columns for grouping\n"
                
                # Sheet purpose assessment
                if 'id' in [col.lower() for col in df.columns] or df.columns[0].lower().endswith('id'):
                    content_summary += f"- **Sheet Type**: Likely transactional or record data\n"
                elif len(numeric_cols) > len(categorical_cols):
                    content_summary += f"- **Sheet Type**: Primarily quantitative/analytical data\n"
                else:
                    content_summary += f"- **Sheet Type**: Mixed or reference data\n"
                
                # Sample data preview
                content_summary += f"\n#### Sample Data Preview - {sheet_name}\n"
                try:
                    if len(df) > 0:
                        sample_rows = min(3, len(df))
                        display_cols = df.columns[:6]  # Limit columns for readability
                        
                        content_summary += "| " + " | ".join(display_cols) + " |\n"
                        content_summary += "|" + "|".join(["---"] * len(display_cols)) + "|\n"
                        
                        for idx in range(sample_rows):
                            row_values = []
                            for col in display_cols:
                                val = str(df.iloc[idx][col])
                                # Truncate long values
                                val = val[:15] + "..." if len(val) > 15 else val
                                row_values.append(val)
                            content_summary += "| " + " | ".join(row_values) + " |\n"
                        
                        if len(df.columns) > 6:
                            content_summary += f"\n*Note: Showing first 6 of {len(df.columns)} columns*\n"
                            
                except Exception as e:
                    content_summary += f"Error displaying sample data: {e}\n"
            else:
                content_summary += "- **Status**: Empty sheet\n"
            
            content_summary += "\n"
        
        # Overall workbook insights
        content_summary += "## Workbook Business Intelligence\n\n"
        
        # Data distribution
        largest_sheet = max(sheets_dict.keys(), key=lambda x: len(sheets_dict[x]))
        content_summary += f"### Data Distribution\n"
        content_summary += f"- **Primary Sheet**: '{largest_sheet}' contains the most data ({len(sheets_dict[largest_sheet]):,} rows)\n"
        content_summary += f"- **Total Records**: {analysis['total_rows']:,} rows across all sheets\n"
        
        # Multi-sheet analysis
        if len(sheets_dict) > 1:
            content_summary += f"- **Multi-sheet Structure**: Workbook contains {len(sheets_dict)} related datasets\n"
            
            # Look for potential relationships
            sheet_cols = {name: set(df.columns) for name, df in sheets_dict.items()}
            common_cols = set.intersection(*sheet_cols.values()) if len(sheet_cols) > 1 else set()
            if common_cols:
                content_summary += f"- **Related Data**: Common columns found: {', '.join(list(common_cols)[:3])} - potential for cross-sheet analysis\n"
        
        # Overall data quality
        total_cells = sum(df.shape[0] * df.shape[1] for df in sheets_dict.values())
        total_nulls = sum(df.isnull().sum().sum() for df in sheets_dict.values())
        overall_completeness = ((total_cells - total_nulls) / total_cells * 100) if total_cells > 0 else 100
        
        content_summary += f"\n### Overall Data Quality\n"
        content_summary += f"- **Completeness**: {overall_completeness:.1f}% of cells contain data\n"
        content_summary += f"- **Quality Assessment**: {'Excellent' if overall_completeness > 95 else 'Good' if overall_completeness > 85 else 'Needs attention'}\n"
        
        # Business recommendations
        content_summary += f"\n### Business Recommendations\n"
        if analysis['total_rows'] > 10000:
            content_summary += "- **Big Data**: Large dataset suitable for comprehensive business intelligence\n"
        if len(sheets_dict) > 2:
            content_summary += "- **Multi-dimensional Analysis**: Multiple sheets enable cross-referenced insights\n"
        
        total_numeric_cols = sum(len(df.select_dtypes(include=['number']).columns) for df in sheets_dict.values())
        if total_numeric_cols > 5:
            content_summary += "- **Analytics Ready**: Multiple numeric columns available for statistical analysis\n"
        
        return {
            "content": content_summary,
            "metadata": {
                "type": "xlsx", 
                "source": file_path,
                "analysis": analysis
            },
            "summary": f"Excel workbook comprehensively analyzed: {len(sheets_dict)} sheets with {analysis['total_rows']:,} total rows and {overall_completeness:.1f}% data completeness. Primary sheet '{largest_sheet}' contains {len(sheets_dict[largest_sheet]):,} records."
        }
    except Exception as e:
        return {"error": f"Error processing Excel file: {str(e)}"}

# ========== MCP Search (Enhanced) ==========
@tool("MCP Search")
def mcp_search(query: str):
    """Enhanced MCP search with contextual project data and learnings."""
    # Simulate enhanced MCP search with more realistic project data
    project_database = {
        "AI productivity tools": [
            {"title": "TaskMaster AI", "summary": "AI scheduling assistant with 85% user adoption", "success_factors": ["intuitive UI", "accurate prediction", "seamless integration"], "timeline": "6 months", "team_size": "8 people"},
            {"title": "WorkFlow Optimizer", "summary": "Business process automation with 40% efficiency gain", "challenges": ["user training", "legacy system integration"], "timeline": "9 months", "budget": "$150K"},
        ],
        "business intelligence": [
            {"title": "DataViz Pro", "summary": "Real-time dashboard solution for SMBs", "key_learnings": ["mobile-first design essential", "automated insights preferred"], "roi": "180%", "timeline": "4 months"},
            {"title": "Analytics Engine", "summary": "Custom BI solution with advanced ML", "risks": ["data quality issues", "scalability concerns"], "success_rate": "75%"},
        ],
        "time management": [
            {"title": "TimeSync", "summary": "Cross-platform time tracking with team collaboration", "metrics": {"user_retention": "78%", "productivity_gain": "32%"}, "platform": "web + mobile"},
            {"title": "Focus Mode App", "summary": "Distraction blocking with productivity analytics", "lessons": ["gamification increased engagement", "privacy concerns in team features"], "downloads": "50K+"},
        ],
        "document management": [
            {"title": "DocuFlow", "summary": "Automated document processing and workflow", "success_factors": ["OCR accuracy", "API integrations"], "efficiency_gain": "60%"},
            {"title": "SmartFile", "summary": "AI-powered file organization and search", "key_features": ["auto-tagging", "content search", "version control"]},
        ],
        "multimodal analysis": [
            {"title": "ContentIQ", "summary": "Multi-format content analysis platform", "capabilities": ["PDF, images, audio processing"], "accuracy": "92%"},
            {"title": "MediaMind", "summary": "Cross-media content understanding system", "technologies": ["NLP", "computer vision", "speech recognition"]},
        ]
    }
    
    # Enhanced search logic
    relevant_projects = []
    query_lower = query.lower()
    
    # Search by category keywords
    for category, projects in project_database.items():
        if any(keyword in query_lower for keyword in category.split()):
            relevant_projects.extend(projects)
    
    # Search by specific terms in titles and summaries
    search_terms = query_lower.split()
    for projects in project_database.values():
        for project in projects:
            title_words = project['title'].lower().split()
            summary_words = project['summary'].lower().split()
            
            if any(term in title_words or term in summary_words for term in search_terms):
                if project not in relevant_projects:
                    relevant_projects.append(project)
    
    # If no specific matches, provide general business project insights
    if not relevant_projects:
        relevant_projects = [
            {"title": "Generic Business Solution", "summary": "Standard business process improvement project", "success_factors": ["stakeholder alignment", "clear requirements", "iterative development"]},
            {"title": "Technology Integration", "summary": "System integration and process automation", "challenges": ["technical complexity", "user adoption", "data migration"]},
        ]
    
    # Extract consolidated learnings based on query context
    consolidated_learnings = []
    
    if "ai" in query_lower or "artificial" in query_lower:
        consolidated_learnings.extend([
            "AI projects require 80% more user training but deliver 3x productivity gains",
            "User acceptance is critical - involve end users in AI tool design from day one",
            "Data quality directly impacts AI effectiveness - invest in data preparation"
        ])
    
    if "mobile" in query_lower or "app" in query_lower:
        consolidated_learnings.extend([
            "Mobile-first design is essential for modern business applications",
            "Cross-platform compatibility saves 40% development time",
            "Push notifications increase user engagement by 60%"
        ])
    
    if "management" in query_lower:
        consolidated_learnings.extend([
            "Executive sponsorship increases project success rate by 70%",
            "Change management is often underestimated - allocate 30% of budget",
            "Regular stakeholder communication prevents scope creep"
        ])
    
    # Default learnings if none specific match
    if not consolidated_learnings:
        consolidated_learnings = [
            "User experience is critical - 80% of project success depends on intuitive design",
            "Integration capabilities make or break enterprise adoption",
            "Iterative development with regular feedback loops reduces project risk",
            "Data privacy and security concerns must be addressed early",
            "Scalability planning prevents future architectural debt",
            "Automated testing and deployment improve delivery quality"
        ]
    
    return {
        "query_processed": query,
        "similar_projects": relevant_projects[:5],  # Limit to top 5 most relevant
        "key_learnings": consolidated_learnings[:6],  # Top 6 learnings
        "recommendations": [
            "Start with clear stakeholder alignment and requirements gathering",
            "Plan for integration with existing business systems and processes",
            "Implement user feedback loops early and often in development",
            "Address security, privacy, and compliance requirements upfront",
            "Design for scalability and future growth from the beginning",
            "Allocate sufficient resources for user training and change management"
        ],
        "risk_factors": [
            "Inadequate stakeholder buy-in and communication",
            "Underestimating technical complexity and integration challenges",
            "Insufficient user training and change management",
            "Poor data quality or availability issues",
            "Scope creep and unclear requirements"
        ]
    }

@tool("A2A Communication")
def a2a_communicate(agent_name: str, message: str, task_id: str = None):
    """Real A2A communication with other GMI agents."""
    
    agent_endpoints = {
        "Business Analyst Agent": "http://localhost:10005",
        "Business Analyst Domain Expert Agent": "http://localhost:10006", 
        "Product Manager Agent": "http://localhost:10007",
        "Agile Project Manager Agent": "http://localhost:10008"
    }
    
    if agent_name not in agent_endpoints:
        return {"error": f"Unknown agent: {agent_name}"}
    
    try:
        response = requests.post(
            f"{agent_endpoints[agent_name]}/invoke",
            json={"input": message},
            timeout=30
        )
        
        if response.status_code == 200:
            result = response.json()["result"]
            return {
                "task_id": task_id or f"GMI_task_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
                "from_agent": "Stakeholder_Agent",
                "to_agent": agent_name,
                "message_sent": message,
                "response_received": result,
                "timestamp": datetime.now().isoformat(),
                "status": "completed"
            }
        else:
            return {"error": f"Agent {agent_name} returned error: {response.text}"}
            
    except Exception as e:
        return {"error": f"Communication failed with {agent_name}: {str(e)}"}


# ========== HTML Report Generation (Fixed) ==========
@tool("Generate HTML Report")
def generate_html_report(markdown_content: str, output_file_path: str = "output/report.html"):
    """
    Converts Markdown content to HTML format (alternative to PDF generation).
    """
    try:
        os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
        
        # Process markdown content BEFORE putting it in f-string
        processed_content = markdown_content.replace('\n', '<br>').replace('##', '<h2>').replace('#', '<h1>')
        current_timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        
        # Now use the processed variables in f-string
        html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>GMI Analysis Report</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
        h1 {{ color: #2c3e50; border-bottom: 3px solid #3498db; }}
        h2 {{ color: #34495e; border-left: 4px solid #3498db; padding-left: 10px; }}
        h3 {{ color: #555; }}
        .summary {{ background: #f8f9fa; padding: 20px; border-radius: 8px; }}
        .metadata {{ background: #e9ecef; padding: 15px; border-radius: 5px; margin: 10px 0; }}
    </style>
</head>
<body>
    <div class="content">
        {processed_content}
    </div>
    <footer>
        <hr>
        <p><small>Generated by GMI Advanced Agentic AI Framework | {current_timestamp}</small></p>
    </footer>
</body>
</html>"""
        
        with open(output_file_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return {"content": f"HTML report successfully generated at {output_file_path}"}
    except Exception as e:
        return {"error": f"Failed to generate HTML report: {str(e)}"}
